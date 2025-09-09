import os
import io
import uuid
import re
from typing import List, Tuple, Optional
import logging


import fitz  # PyMuPDF
import docx
from pptx import Presentation
from dotenv import load_dotenv
from openai import OpenAI

import chromadb
from sentence_transformers import SentenceTransformer

# -------------------------
# Logging setup
# -------------------------
logger = logging.getLogger("studysqueeze.utils")
logging.basicConfig(level=logging.INFO)

# -------------------------
# Environment & Providers
# -------------------------
load_dotenv()

def get_openrouter_client() -> Optional[OpenAI]:
    key = os.getenv("OPENROUTER_API_KEY")
    if not key:
        logger.error("OPENROUTER_API_KEY not found in environment variables")
        return None
    return OpenAI(
        base_url="https://openrouter.ai/api/v1",
        api_key=key,
    )

openrouter_client = get_openrouter_client()
_embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
chroma_client = chromadb.Client()

# -------------------------
# File Parsing
# -------------------------
def extract_text_from_file(filename: str, file_bytes: bytes) -> str:
    """
    Extract plain text from supported files (PDF, DOCX, PPTX).
    """
    filename = filename.lower()
    try:
        if filename.endswith(".pdf"):
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            return "\n".join(page.get_text("text") for page in doc)
        elif filename.endswith(".docx"):
            docx_doc = docx.Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in docx_doc.paragraphs if p.text.strip())
        elif filename.endswith(".pptx"):
            prs = Presentation(io.BytesIO(file_bytes))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n".join(t.strip() for t in texts if t.strip())
        else:
            logger.warning(f"Unsupported file format: {filename}")
            return f"Unsupported file format: {filename}"
    except Exception as e:
        logger.error(f"Error extracting text from file: {e}")
        return f"Error extracting text: {e}"

# -------------------------
# Chunking helpers
# -------------------------
def _normalize_whitespace(text: str) -> str:
    """Collapse multiple whitespaces into single spaces."""
    return re.sub(r"\s+", " ", text).strip()

def chunk_text(text: str, chunk_size: int = 900, overlap: int = 120) -> List[str]:
    """
    Split text into overlapping chunks for embedding.
    """
    text = _normalize_whitespace(text)
    chunks, i, n = [], 0, len(text)
    while i < n:
        chunk = text[i:i + chunk_size]
        chunks.append(chunk)
        i += max(chunk_size - overlap, 1)
    return chunks

# -------------------------
# Vector DB helper
# -------------------------
def build_vector_store(extracted_text: str):
    """
    Create or refresh the 'study_doc' collection in Chroma.
    Returns the Chroma collection object.
    """
    try:
        chroma_client.delete_collection("study_doc")
    except Exception:
        pass
    collection = chroma_client.create_collection(name="study_doc")
    chunks = chunk_text(extracted_text)
    if not chunks:
        logger.error("No text chunks could be generated from the document.")
        raise ValueError("No text chunks could be generated from the document.")
    ids = [str(uuid.uuid4()) for _ in chunks]
    embeddings = _embedding_model.encode(chunks, convert_to_numpy=True)
    collection.add(
        documents=chunks,
        ids=ids,
        embeddings=embeddings.tolist(),
        metadatas=[{"source": "uploaded"} for _ in chunks],
    )
    return collection

def get_vector_store():
    """
    Retrieve the existing Chroma collection.
    """
    try:
        return chroma_client.get_collection(name="study_doc")
    except Exception as e:
        logger.error(f"Error accessing vector store: {e}")
        return None

def retrieve_context(collection, user_query: str, top_k: int = 6) -> List[Tuple[str, float]]:
    """
    Retrieve top-k relevant chunks for a query.
    Returns list of (chunk_text, distance).
    """
    try:
        query_emb = _embedding_model.encode(user_query).tolist()
        res = collection.query(
            query_embeddings=[query_emb],
            n_results=top_k,
            include=["documents", "distances"],
        )
        docs = res.get("documents", [[]])[0]
        dists = res.get("distances", [[]])[0]
        return list(zip(docs, dists))
    except Exception as e:
        logger.error(f"Error retrieving context: {e}")
        return []

# -------------------------
# Prompt scaffolding
# -------------------------
MODE_PROMPTS = {
    "quick-recap": (
        "Summarize the material for last-minute revision in crisp bullet points. "
        "Prioritize definitions, theorems, formulas, dates, and cause→effect relationships. "
        "Keep bullets short (max ~15 words) and group them by subtopic."
    ),
    "analogy": (
        "Explain the core ideas using simple analogies and everyday scenarios. "
        "Draw comparisons that are concrete and intuitive. End each analogy with a one-line takeaway."
    ),
    "exam-cram": (
        "Create an exam cheat sheet: key ideas, must-know formulas, common traps, and FAQs. "
        "Use headings and bullets. Make it scannable and compact."
    ),
    "challenge": (
        "Produce a short summary, then generate 5 challenge Q&A pairs that probe real understanding. "
        "Mix conceptual, edge-case, and application-style questions. Give concise answers."
    ),
    "mnemonics": (
        "Extract the most memorization-heavy facts and create mnemonics, acronyms, or short rhymes. "
        "Show the fact first, then the mnemonic in [Mnemonic: ...]."
    ),
}

SYSTEM_INSTRUCTION = (
    "You are StudySqueeze — an extraordinary AI study assistant that helps students learn fast and remember longer. "
    "Always write with clarity and structure. Prefer short, punchy bullets. Bold key terms. "
    "Use headings, subheadings, bullets, and tables when helpful. "
    "When appropriate, add memory aids (mnemonics/analogies), call out common pitfalls, and highlight exam-likely content. "
    "Be supportive and motivating."
)

def build_rag_prompt(user_prompt: str, retrieved_chunks: List[Tuple[str, float]], mode: str) -> str:
    """
    Build the final RAG prompt combining retrieved context + mode-specific instructions.
    """
    mode_instruction = MODE_PROMPTS.get(mode, MODE_PROMPTS["quick-recap"])
    MAX_CONTEXT_CHARS = 6000
    context_text = ""
    for chunk, _ in retrieved_chunks:
        if len(context_text) + len(chunk) + 2 <= MAX_CONTEXT_CHARS:
            context_text += chunk + "\n\n"
        else:
            break
    return (
        f"{SYSTEM_INSTRUCTION}\n\n"
        f"Mode: {mode}\n"
        f"Mode instructions: {mode_instruction}\n\n"
        f"Use ONLY the following context to answer the user. "
        f"If the answer is not found in context, explicitly say so.\n\n"
        f"Context:\n{context_text.strip()}\n\n"
        f"User request:\n{user_prompt}\n\n"
        f"Respond with clean sections and **bold** keywords."
    )

# -------------------------
# LLM call (Gemini via OpenRouter)
# -------------------------
def query_openrouter_with_rag(user_prompt: str, vector_store, mode: str = "default") -> str:
    """
    Retrieves context from the collection and queries Gemini 2.0 Flash Experimental via OpenRouter.
    """
    if not openrouter_client:
        logger.error("OpenRouter client not initialized.")
        return "Error: OpenRouter client not initialized."
    collection = get_vector_store()
    if not collection:
        return "Error: Could not access vector store."
    retrieved = retrieve_context(collection, user_prompt, top_k=6)
    if not retrieved:
        return "No relevant context found in your document."
    prompt = build_rag_prompt(user_prompt, retrieved, mode)
    try:
        resp = openrouter_client.chat.completions.create(
            model="meta-llama/llama-3.3-70b-instruct:free",
            messages=[
                {"role": "system", "content": SYSTEM_INSTRUCTION},
                {"role": "user", "content": prompt},
            ],
            extra_headers={
                "HTTP-Referer": "http://localhost:3000",
                "X-Title": "StudySqueeze",
            },
        )
        logger.info("OpenRouter response received.")
        content = None
        try:
            if hasattr(resp, "choices") and resp.choices:
                choice = resp.choices[0]
                if hasattr(choice, "message") and hasattr(choice.message, "content"):
                    content = choice.message.content
                elif isinstance(choice, dict) and "message" in choice and "content" in choice["message"]:
                    content = choice["message"]["content"]
        except Exception as e:
            logger.error(f"Error parsing OpenRouter response: {e}")
            return f"Error parsing OpenRouter response: {e}"
        if content:
            return content.strip()
        else:
            logger.error("Error: Could not parse content from OpenRouter response.")
            return "Error: Could not parse content from OpenRouter response."
    except Exception as e:
        logger.error(f"Error while querying OpenRouter: {e}")
        return f"Error while querying OpenRouter: {e}"